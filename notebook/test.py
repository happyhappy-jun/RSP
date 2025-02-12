import glob
import os
import random
import numpy as np
import torch
import torchvision.transforms as T
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from decord import VideoReader, cpu
from torchvision.transforms.functional import InterpolationMode
from transformers import AutoModel, AutoTokenizer
from lmdeploy.vl.constants import IMAGE_TOKEN
from lmdeploy import pipeline, TurbomindEngineConfig
from lmdeploy.vl import load_image
from pptx.enum.text import MSO_AUTO_SIZE

class VideoFrameAnalyzer:
    def __init__(self, model_path='OpenGVLab/InternVL2_5-8B'):
        self.pipe = pipeline(model_path, backend_config=TurbomindEngineConfig(session_len=8192, tp=2))
        
    def analyze_frames(self, frame1, frame2):
        prompt = f"""Image-1: {IMAGE_TOKEN}\nImage-2: {IMAGE_TOKEN}\nYou are a movement analyzer specialized in comparing consecutive video frames. Analyze and describe changes between frames using the following guidelines:
FORMAT:
- Provide one clear sentence summarizing the overall dynamics and activity
- Describe qualitative and relative dynamics between frames
- Use precise directional terms (left, right, up, down, forward, backward)
- Focus on observable, concrete changes

ANALYZE THE FOLLOWING ELEMENTS:
- Main Subject/Object:
    - Position: Track center of mass movement
    - Rotation: Note any turns or spins
    - Orientation: Describe facing direction
    - State Changes: Document visible changes in form, shape, or appearance
- Background:
    - Note any changes in background elements
    - Identify moving vs static elements

Keep descriptions concise, objective, and focused on visible changes between frames."""

        frame1 = Image.fromarray(frame1)
        frame2 = Image.fromarray(frame2)

        response = self.pipe((prompt, [frame1, frame2]))
        print(response.text)
        return response.text

class LLMDemoPowerPointGenerator:
    def __init__(self, output_dir="ppt_demos", max_distance=48):
        self.output_dir = output_dir
        self.media_dir = os.path.join(output_dir, "media")
        self.max_distance = max_distance
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.media_dir, exist_ok=True)
        self.analyzer = VideoFrameAnalyzer()

    def sample_frames(self, video_path):
        vr = VideoReader(video_path, ctx=cpu(0))
        seg_len = len(vr)
        least_frames_num = self.max_distance + 1
        
        if seg_len >= least_frames_num:
            idx_cur = random.randint(0, seg_len - least_frames_num)
            interval = random.randint(4, self.max_distance)
            idx_fut = idx_cur + interval
        else:
            indices = random.sample(range(seg_len), 2)
            indices.sort()
            idx_cur, idx_fut = indices
            
        frame_cur = vr[idx_cur].asnumpy()
        frame_fut = vr[idx_fut].asnumpy()
        
        return frame_cur, frame_fut, idx_cur, idx_fut

    def save_frame(self, frame, filename):
        frame_path = os.path.join(self.media_dir, filename)
        Image.fromarray(frame).save(frame_path)
        return frame_path

    def _add_title_slide(self, prs, title="Video Frame Analysis"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = "Frame-by-Frame Analysis with LLM"
        
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 44, 44)

    def _add_frame_analysis_slides(self, prs, video_path, frame_cur, frame_fut, frame_indices, llm_analysis):
        # Frames slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = f"Frame Analysis (Frames {frame_indices[0]} → {frame_indices[1]})"
        
        # Add frames side by side
        width = Inches(5)
        height = Inches(4)
        
        # Save and add first frame
        frame1_path = self.save_frame(frame_cur, f"frame_{frame_indices[0]}.jpg")
        left = Inches(1)
        top = Inches(1.8)
        slide.shapes.add_picture(frame1_path, left, top, width, height)
        
        # Save and add second frame
        frame2_path = self.save_frame(frame_fut, f"frame_{frame_indices[1]}.jpg")
        left = Inches(7)
        slide.shapes.add_picture(frame2_path, left, top, width, height)
        
        # Add frame numbers
        for i, (left, idx) in enumerate([(Inches(1), frame_indices[0]), (Inches(7), frame_indices[1])]):
            txBox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.3))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            p.text = f"Frame {idx}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(11)
        
        # Analysis slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = "LLM Analysis"
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(10)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Clear any existing paragraphs
        for _ in range(len(tf.paragraphs)):
            p = tf.paragraphs[0]
            tr = p._element
            tr.getparent().remove(tr)
        
        # Add text paragraph by paragraph
        for paragraph_text in llm_analysis.split('\n'):
            p = tf.add_paragraph()
            p.text = paragraph_text
            p.font.size = Pt(12)
            p.font.name = 'Calibri'
            p.space_after = Pt(6)
        
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def generate_presentation(self, videos, title="Video Frame Analysis"):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        self._add_title_slide(prs, title)
        
        for video_path in videos:
            frame_cur, frame_fut, idx_cur, idx_fut = self.sample_frames(video_path)
            llm_analysis = self.analyzer.analyze_frames(frame_cur, frame_fut)
            self._add_frame_analysis_slides(prs, video_path, frame_cur, frame_fut, 
                                         (idx_cur, idx_fut), llm_analysis)
            
        output_path = os.path.join(self.output_dir, "video_analysis.pptx")
        prs.save(output_path)
        return output_path


# Example usage
if __name__ == "__main__":
    videos = glob.glob("/data/RSP/random_videos_to_send/*.mp4")
    
    generator = LLMDemoPowerPointGenerator(max_distance=48)
    ppt_path = generator.generate_presentation(videos, "Robot Task Analysis")
    print(f"Presentation generated at: {ppt_path}")